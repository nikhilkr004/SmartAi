from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_study_ppt(notes: str, topic: str = "Class Session", output_path: str = "output.pptx") -> str:
    """
    Generates a world-class PowerPoint presentation from AI notes.
    """
    prs = Presentation()

    # 1. TITLE SLIDE
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = topic.title()
    subtitle.text = "AI-Generated Smart Study Guide\nFocus: Mastery & Insight"

    # 2. SECTIONS
    lines = notes.split("\n")
    current_slide = None
    current_content = []

    def flush_slide():
        nonlocal current_slide, current_content
        if current_slide and current_content:
            text_frame = current_slide.placeholders[1].text_frame
            text_frame.word_wrap = True
            for point in current_content[:6]: # Limit bullet points for legibility
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
            current_content = []

    for line in lines:
        line = line.strip()
        if not line: continue

        if line.startswith("# "):
            flush_slide()
            layout = prs.slide_layouts[1] # Title and Content
            current_slide = prs.slides.add_slide(layout)
            current_slide.shapes.title.text = line[2:].strip()
            
        elif line.startswith("## "):
            # If we have a subtopic, maybe a new slide? 
            # Or just add a bold line in content?
            current_content.append(f"--- {line[3:].strip()} ---")
            
        elif line.startswith("- ") or line.startswith("* "):
            current_content.append(line[2:].strip())
        elif line[0].isdigit() and line[1] == ".":
            current_content.append(line[3:].strip())
        else:
            # Body text
            if len(line) > 10:
                current_content.append(line)

    flush_slide()

    # 3. CLOSING SLIDE
    closing_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(closing_layout)
    slide.shapes.title.text = "Key Takeaways"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Review your masterclass cheat sheet.\n• Practice the visual logic flows.\n• Mastery comes from repetition."

    prs.save(output_path)
    return output_path
